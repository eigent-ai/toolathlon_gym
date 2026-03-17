"""
Evaluation for canvas-grade-report-ppt-email task.

Checks:
1. PowerPoint Grade_Report.pptx with correct grade data
2. Email sent to instructors@university.edu ( DB check)
"""
import argparse
import sys
import os
from pathlib import Path

PASS_COUNT = 0
FAIL_COUNT = 0

COURSE_ID = 1


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        d = (detail[:300]) if len(detail) > 300 else detail
        print(f"  [FAIL] {name}: {d}")


def get_expected_data():
    """Query PostgreSQL for expected course grade data."""
    import psycopg2

    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()

    cur.execute("SELECT name FROM canvas.courses WHERE id = %s", (COURSE_ID,))
    course_name = cur.fetchone()[0]

    cur.execute("""
        SELECT
            grades->>'current_grade' as grade,
            COUNT(*) as cnt,
            ROUND(AVG((grades->>'current_score')::float)::numeric, 2) as avg_score
        FROM canvas.enrollments
        WHERE course_id = %s AND grades->>'current_grade' IS NOT NULL
        GROUP BY grades->>'current_grade'
        ORDER BY grades->>'current_grade'
    """, (COURSE_ID,))
    grade_dist = [(r[0], int(r[1]), float(r[2])) for r in cur.fetchall()]

    cur.execute("""
        SELECT COUNT(*) as total_enrolled,
               SUM(CASE WHEN grades->>'current_grade' IS NOT NULL THEN 1 ELSE 0 END) as graded,
               ROUND(AVG(CASE WHEN grades->>'current_score' IS NOT NULL
                   THEN (grades->>'current_score')::float END)::numeric, 2) as overall_avg
        FROM canvas.enrollments
        WHERE course_id = %s
    """, (COURSE_ID,))
    row = cur.fetchone()
    summary = (int(row[0]), int(row[1]), float(row[2]))

    conn.close()
    return course_name, grade_dist, summary


def check_pptx(workspace, course_name, grade_dist, summary):
    """Check Grade_Report.pptx for correctness."""
    from pptx import Presentation

    print("\n=== Checking PowerPoint ===")
    pptx_path = Path(workspace) / "Grade_Report.pptx"

    if not pptx_path.exists():
        check("PPTX file exists", False, f"Not found: {pptx_path}")
        return
    check("PPTX file exists", True)

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    check("Has at least 4 slides", len(slides) >= 4, f"Got {len(slides)}")

    # Collect all text
    all_text = []
    for slide in slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_texts.append(paragraph.text)
        all_text.append("\n".join(slide_texts))

    full_text = "\n".join(all_text).lower()

    # Check title slide
    if len(all_text) > 0:
        check("Title slide has 'Grade Analysis'",
              "grade analysis" in all_text[0].lower(),
              f"Text: {all_text[0][:100]}")

    # Check grade distribution data appears
    for grade, cnt, avg_s in grade_dist:
        check(f"Grade '{grade}' count {cnt} in presentation",
              str(cnt) in full_text,
              f"'{cnt}' not found")

    # Check summary values
    total_enrolled, graded, overall_avg = summary
    check("Total enrolled appears",
          str(total_enrolled) in full_text,
          f"Expected {total_enrolled}")

    check("Overall avg score appears",
          str(overall_avg) in full_text,
          f"Expected {overall_avg}")

    # Check pass rate
    total_graded = sum(cnt for _, cnt, _ in grade_dist)
    pass_count = sum(cnt for g, cnt, _ in grade_dist if g != 'Fail')
    pass_rate = round(pass_count / total_graded * 100, 1) if total_graded > 0 else 0
    check("Pass rate appears",
          str(pass_rate) in full_text,
          f"Expected {pass_rate}")


def check_email(course_name, summary):
    """Check email was sent ( DB check)."""
    import psycopg2

    print("\n=== Checking Email ===")
    try:
        conn = psycopg2.connect(
            host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
            user="eigent", password="camel"
        )
        cur = conn.cursor()
    except Exception as e:
        check("DB connection", False, str(e))
        return

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_emails = cur.fetchall()
    cur.close()
    conn.close()

    target = "instructors@university.edu"
    found = None
    for subj, from_addr, to_addr, body in all_emails:
        to_str = str(to_addr or "").lower()
        if target in to_str:
            found = (subj, from_addr, to_addr, body)
            break

    check(f"Email sent to {target}", found is not None,
          f"Found {len(all_emails)} total emails")

    if found:
        subj, _, _, body = found
        body_lower = (body or "").lower()
        subj_lower = (subj or "").lower()

        check("Email subject mentions grade/analytics",
              "grade" in subj_lower or "analytics" in subj_lower or "analysis" in subj_lower,
              f"Subject: {(subj or '')[:100]}")

        check("Email body mentions overall avg score",
              str(summary[2]) in (body or ""),
              f"Expected {summary[2]} in body")

        total_graded_count = summary[1]
        total_enrolled = summary[0]
        check("Email body mentions total enrolled",
              str(total_enrolled) in (body or ""),
              f"Expected {total_enrolled} in body")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("Fetching expected data...")
    course_name, grade_dist, summary = get_expected_data()
    print(f"  Course: {course_name}")
    print(f"  Grade categories: {len(grade_dist)}")
    print(f"  Summary: enrolled={summary[0]}, graded={summary[1]}, avg={summary[2]}")

    check_pptx(args.agent_workspace, course_name, grade_dist, summary)
    check_email(course_name, summary)

    all_ok = FAIL_COUNT == 0

    print(f"\n=== SUMMARY ===")
    print(f"  Total checks - Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}")
    print(f"  Overall: {'PASS' if all_ok else 'FAIL'}")

    if all_ok:
        print("\nPass all tests!")
        sys.exit(0)
    else:
        print("\nSome checks failed.")
        sys.exit(1)
